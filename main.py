import pandas as pd
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches

# 1. Data Preparation
data = {
    'Industry': ['Bollywood', 'Tollywood', 'Kollywood'],
    '2000-2012_Avg_Revenue': [450, 200, 180],
    '2013-2026_Avg_Revenue': [550, 800, 650]
}
df = pd.DataFrame(data)

# 2. Create Bar Chart
plt.figure(figsize=(10, 6))
df.plot(x='Industry', y=['2000-2012_Avg_Revenue', '2013-2026_Avg_Revenue'], kind='bar')
plt.title('Box Office Performance Comparison (in Crores)')
plt.ylabel('Revenue')
plt.savefig('chart.png')

# 3. Create PowerPoint
prs = Presentation()

def add_slide(title, content):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    slide.placeholders[1].text = content

# Slide 1: Title
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "The Shift: Bollywood vs South Indian Cinema"
slide.placeholders[1].text = "2000 - 2026 Analysis"

# Add slides 2-9 content
content_map = {
    "The Rise of Regional Cinema": "How Tollywood and Kollywood shifted from regional to Pan-India.",
    "Bollywood's Stagnation": "Analysis of repetitive scripts and lack of cultural connection.",
    "Box Office Trends": "See attached chart.",
    "Top 5 Bollywood Stars": "SRK, Salman Khan, Aamir Khan, Ranbir Kapoor, Hrithik Roshan.",
    "Top 5 Tollywood Stars": "Prabhas, Allu Arjun, Jr NTR, Ram Charan, Mahesh Babu.",
    "Top 5 Kollywood Stars": "Rajinikanth, Kamal Haasan, Vijay, Ajith, Vikram.",
    "Quality Shift (2000-2025)": "Shift from masala to high-concept storytelling in the South.",
    "Conclusion": "The future of Indian Cinema lies in regional innovation."
}

for title, text in content_map.items():
    add_slide(title, text)

# Add chart to slide 4
slide = prs.slides[3]
slide.shapes.add_picture('chart.png', Inches(1), Inches(2), width=Inches(6))

prs.save('Indian_Cinema_Analysis.pptx')
print("Presentation created successfully!")
